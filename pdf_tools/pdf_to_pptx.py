import pptx # pip install python-pptx
import cv2

def pdfToImages( strPdfFilename ):
    """
    Take a pdf file and export each files to a folder.
    Return the list of generated images files
    """
        
    import fitz # pip install PyMuPDF # https://pypi.org/project/PyMuPDF/#files
    
    doc = fitz.open(strPdfFilename)
    print("INF: pdfToImages: doc page number: %d" % (doc.pageCount) )
    print("INF: pdfToImages: doc meta: %s" % (doc.metadata) )
    
    listImgs = []
    for nNumPage in range(doc.pageCount):
        print( "INF: pdfToImages: generating for page %d" % (nNumPage+1) )
        page = doc.loadPage(nNumPage) #number of page
        if 1:
            zoom_x = 2.0  # 2 => double resolution 72 dpi => 144
            zoom_y = 2.0  # vertical zoom
            mat = fitz.Matrix(zoom_x, zoom_y)  # zoom factor 2 in each dimension
            pix = page.getPixmap(matrix = mat)
        else:
            pix = page.getPixmap()
        output = "/tmp/outfile%04d.png" % nNumPage
        pix.writePNG(output)
        listImgs.append(output)

    doc.close()
    
    return listImgs

def imagesToPres( listFileImages, strDestFilenamePPT = "generated.pptx" ):
    """
    Take a list of images and paste them in a ppt
    """
    pres = pptx.Presentation()
    blank_slide_layout = pres.slide_layouts[6]


    multiplier = 9525

    wmax,hmax = 1920,1080 # HD ratio
    wmax,hmax = 2100,2970 # A4

    if 0:
        # find maximum, based on images (interesing if many images have equivalent size)
        # else we will remains on a classic ratio
        for nNumPage, slideimgfilename in enumerate(listFileImages):
            slideimg = cv2.imread(slideimgfilename)
            height,width = slideimg.shape[:2]

            if width > wmax:
                wmax = width
                
            if height > hmax:
                hmax = height


    # Set slide dimensions
    pres.slide_width = wmax*multiplier
    pres.slide_height = hmax*multiplier

    for nNumPage, slideimgfilename in enumerate(listFileImages):
        slideimg = cv2.imread(slideimgfilename)
        height,width = slideimg.shape[:2]
        
        fit_w_h = int(wmax * height/width) # h when fitting w
        fit_h_w = int(hmax * width/height)
        # Cas 1: NoLose: on ne veut pas perdre des bouts d'images quitte a avoir du blanc sur les cotes
        # Cas 2: NoWhite: On ne veut pas avoir de blanc quitte a perdre un peu d'images
        bNoLose = 0 # else it's NoWhite
        if (bNoLose and fit_w_h > hmax) or (not bNoLose and fit_w_h < hmax):
            # on garde fit_h_w
            wnew, hnew = fit_h_w,hmax
        else:
            wnew, hnew = wmax,fit_w_h
        
        

        # Add slide
        slide = pres.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(slideimgfilename, multiplier*(wmax-wnew)/2, multiplier*(hmax-hnew)/2, width=wnew*multiplier, height=hnew*multiplier)

    base_name = "generated"
    print("INF: imagesToPres: Saving to file: " + strDestFilenamePPT )
    pres.save(strDestFilenamePPT)
# imagesToPres - end

def convertPdfToPpt( strSrcPdfFilename ):
    strOut = strSrcPdfFilename.replace(".pdf",".pptx")
    li = pdfToImages(strSrcPdfFilename)
    imagesToPres(li,strOut)

    
def auto_test_im_to_pptx():
    # imgs to pptx
    aFiles = [
                    "../data/fruit_face.jpg",
                    "../data/face_bw5.jpg",
                    "../data/inconnus.jpg",
                ]
                
    imagesToPres(aFiles)
    
def auto_test():
    li = pdfToImages("devoirs.pdf")
    imagesToPres(li)
    
if __name__ == "__main__":
    #~ auto_test_im_to_pptx()
    #~ auto_test()
    convertPdfToPpt("nda.pdf")
    
